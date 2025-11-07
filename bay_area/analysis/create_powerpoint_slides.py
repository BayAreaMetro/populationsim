#!/usr/bin/env python3
"""
PowerPoint Slides Generator for Marginal Controls
Creates editable PowerPoint presentation from marginal controls data
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os
from pathlib import Path

class PowerPointGenerator:
    def __init__(self):
        self.output_dir = Path("docs/visualizations/marginal_controls")
        self.output_dir.mkdir(parents=True, exist_ok=True)
        
        # Human-readable labels
        self.income_labels = {
            'hhinc_1': 'Under $30,000',
            'hhinc_2': '$30,000 to $60,000', 
            'hhinc_3': '$60,000 to $100,000',
            'hhinc_4': '$100,000 to $150,000',
            'hhinc_5': '$150,000 to $200,000',
            'hhinc_6': '$200,000 to $250,000',
            'hhinc_7': '$250,000 to $300,000',
            'hhinc_8': '$300,000 and above'
        }
        
        self.age_labels = {
            'pers_age_1': 'Ages 0-17',
            'pers_age_2': 'Ages 18-64', 
            'pers_age_3': 'Ages 65-79',
            'pers_age_4': 'Ages 80+'
        }
        
        self.size_labels = {
            'hh_size_1': '1-person household',
            'hh_size_2': '2-person household',
            'hh_size_3': '3-person household', 
            'hh_size_4': '4-person household',
            'hh_size_5': '5-person household',
            'hh_size_6': '6+ person household'
        }
        
        self.worker_labels = {
            'hh_workers_0': '0 workers',
            'hh_workers_1': '1 worker',
            'hh_workers_2': '2 workers', 
            'hh_workers_3': '3+ workers'
        }
        
        self.occupation_labels = {
            'pers_occ_1': 'Professional and Related',
            'pers_occ_2': 'Management and Business',
            'pers_occ_3': 'Service Occupations',
            'pers_occ_4': 'Sales and Office',
            'pers_occ_5': 'Manual and Military'
        }
        
        self.county_names = {
            1: 'San Francisco County',
            2: 'San Mateo County', 
            3: 'Santa Clara County',
            4: 'Alameda County',
            5: 'Contra Costa County',
            6: 'Solano County',
            7: 'Napa County',
            8: 'Sonoma County',
            9: 'Marin County'
        }

    def create_presentation(self):
        """Create the complete PowerPoint presentation."""
        print("Creating PowerPoint presentation...")
        
        # Create presentation
        prs = Presentation()
        
        # Set slide dimensions (16:9 aspect ratio)
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # Create slides
        self._create_maz_slide(prs)
        self._create_taz_slide(prs) 
        self._create_county_slide(prs)
        
        # Save presentation
        output_path = self.output_dir / "marginal_controls_presentation.pptx"
        prs.save(output_path)
        print(f"✓ PowerPoint presentation saved: {output_path}")
        
        return output_path

    def _create_maz_slide(self, prs):
        """Create MAZ controls slide."""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(1))
        title_frame = title_box.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = "MAZ (Micro Analysis Zone) Controls"
        title_p.font.size = Pt(36)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(44, 62, 80)
        title_p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.6))
        subtitle_frame = subtitle_box.text_frame
        subtitle_p = subtitle_frame.paragraphs[0]
        subtitle_p.text = "Finest Geographic Level: 39,586 zones • Housing Units & Group Quarters • 2020/2010 Interpolation"
        subtitle_p.font.size = Pt(18)
        subtitle_p.font.italic = True
        subtitle_p.font.color.rgb = RGBColor(127, 140, 141)
        subtitle_p.alignment = PP_ALIGN.CENTER
        
        # Left column - Categories
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(6), Inches(4.8))
        left_frame = left_box.text_frame
        left_frame.word_wrap = True
        
        # Add background shape
        left_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(2.1), Inches(6.2), Inches(5))
        left_bg.fill.solid()
        left_bg.fill.fore_color.rgb = RGBColor(173, 216, 230)  # Light blue
        left_bg.line.color.rgb = RGBColor(52, 152, 219)
        left_bg.line.width = Pt(3)
        
        # Move text box to front
        slide.shapes._spTree.remove(left_box._element)
        slide.shapes._spTree.append(left_box._element)
        
        categories_text = """Housing Unit Categories:

• Total Households (numhh_gq)
• University Group Quarters (gq_type_univ)
• Non-institutional Group Quarters (gq_type_noninst)

Geographic Coverage:
• 39,586 MAZ zones across Bay Area
• Finest level of geographic detail
• Building block for TAZ aggregation"""
        
        left_p = left_frame.paragraphs[0]
        left_p.text = categories_text
        left_p.font.size = Pt(16)
        left_p.font.color.rgb = RGBColor(52, 73, 94)
        
        # Right column - Data Sources
        right_box = slide.shapes.add_textbox(Inches(7), Inches(2.2), Inches(6), Inches(4.8))
        right_frame = right_box.text_frame
        right_frame.word_wrap = True
        
        # Add background shape
        right_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.9), Inches(2.1), Inches(6.2), Inches(5))
        right_bg.fill.solid()
        right_bg.fill.fore_color.rgb = RGBColor(211, 211, 211)  # Light gray
        right_bg.line.color.rgb = RGBColor(149, 165, 166)
        right_bg.line.width = Pt(3)
        
        # Move text box to front
        slide.shapes._spTree.remove(right_box._element)
        slide.shapes._spTree.append(right_box._element)
        
        sources_text = """Data Sources & Methodology:

Primary Sources:
• 2020 Census (block level counts)
• 2020/2010 NHGIS interpolation
• Geographic allocation methods

Processing Steps:
• Block-to-MAZ aggregation
• Group quarters identification
• Institutional vs. non-institutional classification
• Quality validation against regional totals"""
        
        right_p = right_frame.paragraphs[0]
        right_p.text = sources_text
        right_p.font.size = Pt(16)
        right_p.font.color.rgb = RGBColor(52, 73, 94)

    def _create_taz_slide(self, prs):
        """Create TAZ controls slide."""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.8))
        title_frame = title_box.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = "TAZ (Traffic Analysis Zone) Controls"
        title_p.font.size = Pt(32)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(44, 62, 80)
        title_p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(12.3), Inches(0.5))
        subtitle_frame = subtitle_box.text_frame
        subtitle_p = subtitle_frame.paragraphs[0]
        subtitle_p.text = "Transportation Level: 4,734 zones • Household & Person Demographics • 2020/2010 Interpolation"
        subtitle_p.font.size = Pt(16)
        subtitle_p.font.italic = True
        subtitle_p.font.color.rgb = RGBColor(127, 140, 141)
        subtitle_p.alignment = PP_ALIGN.CENTER
        
        # Create 3x2 grid
        colors = [
            RGBColor(144, 238, 144),  # Light green
            RGBColor(255, 255, 224),  # Light yellow  
            RGBColor(240, 128, 128),  # Light coral
            RGBColor(176, 196, 222),  # Light steel blue
            RGBColor(245, 222, 179),  # Wheat
            RGBColor(211, 211, 211)   # Light gray
        ]
        
        box_width = Inches(4.2)
        box_height = Inches(2.8)
        
        # Top row
        self._add_taz_box(slide, Inches(0.3), Inches(1.6), box_width, box_height, 
                         "Income Distribution", self._get_income_text(), colors[0])
        self._add_taz_box(slide, Inches(4.6), Inches(1.6), box_width, box_height,
                         "Age Distribution", self._get_age_text(), colors[1])
        self._add_taz_box(slide, Inches(8.9), Inches(1.6), box_width, box_height,
                         "Household Size", self._get_size_text(), colors[2])
        
        # Bottom row  
        self._add_taz_box(slide, Inches(0.3), Inches(4.5), box_width, box_height,
                         "Worker Distribution", self._get_worker_text(), colors[3])
        self._add_taz_box(slide, Inches(4.6), Inches(4.5), box_width, box_height,
                         "Additional Controls", self._get_additional_text(), colors[4])
        self._add_taz_box(slide, Inches(8.9), Inches(4.5), box_width, box_height,
                         "Data Sources", self._get_taz_sources_text(), colors[5])

    def _add_taz_box(self, slide, left, top, width, height, title, content, bg_color):
        """Add a formatted text box for TAZ slide."""
        # Background shape
        bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = bg_color
        bg_shape.line.color.rgb = RGBColor(100, 100, 100)
        bg_shape.line.width = Pt(2)
        
        # Text box
        text_box = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.1), 
                                          width - Inches(0.2), height - Inches(0.2))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        # Title
        title_p = text_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(14)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(44, 62, 80)
        title_p.alignment = PP_ALIGN.CENTER
        
        # Content
        content_p = text_frame.add_paragraph()
        content_p.text = content
        content_p.font.size = Pt(11)
        content_p.font.color.rgb = RGBColor(52, 73, 94)

    def _get_income_text(self):
        """Get income categories text."""
        categories = [f"• {label}" for label in self.income_labels.values()]
        return "Households by Income:\n\n" + "\n".join(categories)

    def _get_age_text(self):
        """Get age categories text.""" 
        categories = [f"• {label}" for label in self.age_labels.values()]
        return "Persons by Age Group:\n\n" + "\n".join(categories)

    def _get_size_text(self):
        """Get household size categories text."""
        categories = [f"• {label}" for label in self.size_labels.values()]
        return "Households by Size:\n\n" + "\n".join(categories)

    def _get_worker_text(self):
        """Get worker categories text."""
        categories = [f"• {label}" for label in self.worker_labels.values()]
        return "Households by Workers:\n\n" + "\n".join(categories)

    def _get_additional_text(self):
        """Get additional categories text."""
        return """Additional Categories:
• Households with Children (Yes/No)
• Single-Person Group Quarters

Geographic Coverage:
• 4,734 TAZ zones
• Transportation modeling units"""

    def _get_taz_sources_text(self):
        """Get TAZ data sources text."""
        return """Data Sources:
• 2020/2010 NHGIS interpolation
• Census block group mapping
• Geographic allocation methods

Processing:
• Spatial disaggregation
• Control harmonization
• Quality validation"""

    def _create_county_slide(self, prs):
        """Create County controls slide."""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8))
        title_frame = title_box.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = "County Level Controls"
        title_p.font.size = Pt(36)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(44, 62, 80)
        title_p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.5))
        subtitle_frame = subtitle_box.text_frame
        subtitle_p = subtitle_frame.paragraphs[0]
        subtitle_p.text = "Regional Level: 9 counties • Employment by Occupation • 2020/2010 Interpolation + SOCP codes"
        subtitle_p.font.size = Pt(16)
        subtitle_p.font.italic = True
        subtitle_p.font.color.rgb = RGBColor(127, 140, 141)
        subtitle_p.alignment = PP_ALIGN.CENTER
        
        # Create 2x2 grid
        colors = [
            RGBColor(255, 182, 193),  # Light pink
            RGBColor(230, 230, 250),  # Lavender
            RGBColor(211, 211, 211),  # Light gray
            RGBColor(224, 255, 255)   # Light cyan
        ]
        
        box_width = Inches(6.2)
        box_height = Inches(2.8)
        
        # Top row
        self._add_county_box(slide, Inches(0.4), Inches(1.8), box_width, box_height,
                           "Geographic Coverage", self._get_counties_text(), colors[0])
        self._add_county_box(slide, Inches(6.9), Inches(1.8), box_width, box_height,
                           "Employment Categories", self._get_occupation_text(), colors[1])
        
        # Bottom row
        self._add_county_box(slide, Inches(0.4), Inches(4.7), box_width, box_height,
                           "Data Sources & Geography", self._get_county_sources_text(), colors[2])
        self._add_county_box(slide, Inches(6.9), Inches(4.7), box_width, box_height,
                           "Processing & Validation", self._get_processing_text(), colors[3])

    def _add_county_box(self, slide, left, top, width, height, title, content, bg_color):
        """Add a formatted text box for County slide."""
        # Background shape
        bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = bg_color
        bg_shape.line.color.rgb = RGBColor(100, 100, 100)
        bg_shape.line.width = Pt(2)
        
        # Text box
        text_box = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.15),
                                          width - Inches(0.3), height - Inches(0.3))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        # Title
        title_p = text_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(16)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(44, 62, 80)
        title_p.alignment = PP_ALIGN.CENTER
        
        # Content
        content_p = text_frame.add_paragraph()
        content_p.text = content
        content_p.font.size = Pt(13)
        content_p.font.color.rgb = RGBColor(52, 73, 94)

    def _get_counties_text(self):
        """Get counties list text."""
        counties = [f"• {name}" for name in self.county_names.values()]
        return "Bay Area Counties:\n\n" + "\n".join(counties)

    def _get_occupation_text(self):
        """Get occupation categories text."""
        categories = [f"• {label}" for label in self.occupation_labels.values()]
        return "Persons by Occupation:\n\n" + "\n".join(categories)

    def _get_county_sources_text(self):
        """Get county data sources text."""
        return """Data Sources:
• 2020/2010 NHGIS interpolation
• SOCP occupation codes
• County-level aggregation
• Census geographic boundaries

Geographic Framework:
• 9-county Bay Area region
• Standard metropolitan area
• Regional planning boundaries"""

    def _get_processing_text(self):
        """Get processing methodology text."""
        return """Employment Classification:
• SOCP to 5-category mapping
• Professional vs. service workers
• Management hierarchy
• Manual labor & military

Quality Assurance:
• Regional employment totals
• Industry sector validation
• Cross-county consistency
• Labor force participation rates"""


def main():
    """Main function to generate PowerPoint presentation."""
    print("=" * 60)
    print("POWERPOINT PRESENTATION GENERATOR")
    print("=" * 60)
    
    generator = PowerPointGenerator()
    output_path = generator.create_presentation()
    
    print("=" * 60)
    print("POWERPOINT GENERATION COMPLETE")
    print("=" * 60)
    print(f"Output file: {output_path}")
    print("\n🎉 PowerPoint presentation created successfully!")
    print("📊 Open the .pptx file in Microsoft PowerPoint to edit")


if __name__ == "__main__":
    main()


